# Python-Pptx Module
from pptx import Presentation
from pptx.util import Inches

def set_paragraphs(slide, paragraphs):
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    text_frame = body_shape.text_frame

    for level, text in paragraphs:
        p = text_frame.add_paragraph()
        p.level = level
        p.text = text

    return slide

def gen_pptx(data, DESIGN_TEMPLATE): 
# Make Presentation - Title
    presentation = Presentation(DESIGN_TEMPLATE)
    # presentation.slide_width = Inches( 16 ) 
    # presentation.slide_height = Inches( 9 )
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1] 
    title.text = data["title_name"]

    # Make Contents Page of Presentation
    contents_slide_layout = presentation.slide_layouts[1]
    contents_slide = presentation.slides.add_slide(contents_slide_layout)
    contents_title = contents_slide.shapes.title
    contents_title.text = "Contents"
    contents_list = contents_slide.placeholders[1]
    paragraphs = []
    for section in data["section_content"]:
        title = ""
        for name in section["section_name"]:
            title += name
        paragraphs.append((1, title))

        # 箇条書きのレイアウトを指定
    set_paragraphs(contents_slide, paragraphs)

    # Generate Body text for each section and insert it on the page.
    ## Chains (LangChain Function) is used to generate the page body text 
    ## with the contents of each section/page being dependent on each other.
    for section in data["section_content"]:
        content_slide_layout = presentation.slide_layouts[1] 
        content_slide = presentation.slides.add_slide(content_slide_layout)
        title = ""
        for name in section["section_name"]:
            title += name
        content_title = content_slide.shapes.title
        content_title.text = title
        
        paragraphs = []
        for content in section["section_context"]:
            paragraphs.append((1, content))

            # 箇条書きのレイアウトを指定
            set_paragraphs(content_slide, paragraphs)


    presentation.save(data["title_name"] + ".pptx")